"""
Генератор Кредитных Отчетов (MVP) — Hugging Face Spaces App
Two tariffs: "Стандарт" (NSD API + LLM commentary) and "Премиум" (full LLM analysis).
"""

import os, json, shutil, tempfile, re
import gradio as gr
from openai import OpenAI
import requests as _requests
from bs4 import BeautifulSoup as _BS

# ---------------------------------------------------------------------------
# LLM setup — OpenRouter (works from Russian IPs, accesses Gemini/Claude/Llama)
# ---------------------------------------------------------------------------

_LLM_MODEL = "google/gemini-2.5-flash"  # or "anthropic/claude-sonnet-4", "meta-llama/llama-3.1-70b-instruct"

def _get_api_keys(user_api_key: str = "") -> list[str]:
    keys = []
    for env_var in ["OPENROUTER_KEY"]:
        k = os.environ.get(env_var, "").strip()
        if k: keys.append(k)
    if user_api_key.strip(): keys.append(user_api_key.strip())
    return keys

def get_gemini_model(user_api_key: str = ""):
    keys = _get_api_keys(user_api_key)
    if not keys:
        raise ValueError("OpenRouter API key not found. Set OPENROUTER_KEY in Secrets or paste in UI.")
    return keys[0]

# ---------------------------------------------------------------------------
# PDF → text extraction (with OCR fallback)
# ---------------------------------------------------------------------------

def extract_text_from_pdf(pdf_path: str) -> str:
    import fitz
    try:
        import pytesseract; from PIL import Image; import io as _io; ocr_available = True
    except ImportError:
        ocr_available = False
    doc = fitz.open(pdf_path)
    pages = []; prev_report_year = None; ocr_count = 0
    for i, page in enumerate(doc):
        text = page.get_text()
        if len(text.strip()) < 100 and ocr_available:
            try:
                print(f"  OCR стр.{i+1}/{len(doc)}...", flush=True)
                pix = page.get_pixmap(dpi=100, colorspace=fitz.csGRAY)
                img = Image.open(_io.BytesIO(pix.tobytes("png")))
                img = img.point(lambda x: 0 if x < 180 else 255, '1')
                text = pytesseract.image_to_string(img, lang='rus', config='--psm 6', timeout=245)
                ocr_count += 1; del pix, img
            except RuntimeError: print(f"  OCR timeout стр.{i+1}", flush=True)
            except Exception as e: print(f"  OCR ошибка стр.{i+1}: {e}", flush=True)
        if text.strip():
            ym = re.findall(r'(?:на\s+31\s+декабря|за)\s+(20[12]\d)\s+год', text, re.IGNORECASE)
            if ym:
                ry = max(ym)
                if ry != prev_report_year:
                    pages.append(f"\n{'='*60}\n=== НАЧАЛО ГОДОВОГО ОТЧЁТА ЗА {ry} ГОД ===\n{'='*60}")
                    prev_report_year = ry
            pages.append(f"--- Страница {i+1} ---\n{text}")
    doc.close()
    if ocr_count > 0: pages.insert(0, f"[OCR: {ocr_count} стр. распознано]")
    return "\n\n".join(pages)

def extract_texts(file_list) -> list[str]:
    texts = []
    if not file_list: return texts
    for f in file_list:
        path = f.name if hasattr(f, "name") else f
        fname = os.path.basename(path)
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == ".pdf": t = extract_text_from_pdf(path)
            else:
                with open(path, "r", encoding="utf-8", errors="ignore") as fh: t = fh.read()
            if t.strip(): texts.append(f"=== НАЧАЛО ДОКУМЕНТА: {fname} ===\n{t}\n=== КОНЕЦ ДОКУМЕНТА: {fname} ===")
        except Exception as e: texts.append(f"[Ошибка: {fname}: {e}]")
    return texts

# ---------------------------------------------------------------------------
# MOEX NSD Client (for "Стандарт" tariff)
# ---------------------------------------------------------------------------

class MOEXClient:
    def __init__(self):
        self.session = _requests.Session()
        self.session.headers.update({"User-Agent": "Mozilla/5.0", "Accept": "application/json"})
        self.authenticated = False
    def login(self) -> bool:
        creds = os.environ.get("NSD", "").strip()
        if not creds or ":" not in creds: print("❌ NSD secret missing", flush=True); return False
        username, password = creds.split(":", 1)
        try:
            r = self.session.get("https://passport.moex.com/login")
            soup = _BS(r.text, 'html.parser'); form = soup.find('form')
            if not form: return False
            ld = {i.get('name'): i.get('value','') for i in form.find_all('input') if i.get('name')}
            for i in form.find_all('input', {'type':'text'}): ld[i.get('name','')] = username
            for i in form.find_all('input', {'type':'password'}): ld[i.get('name','')] = password
            action = form.get('action', 'https://passport.moex.com/login')
            if action.startswith('/'): action = f"https://passport.moex.com{action}"
            r = self.session.post(action, data=ld, allow_redirects=True)
            self.authenticated = r.status_code == 200 and 'authenticate' not in r.url.lower()
            print(f"{'✅' if self.authenticated else '❌'} MOEX login", flush=True)
            return self.authenticated
        except Exception as e: print(f"❌ MOEX error: {e}", flush=True); return False
    def get_report_list(self, inn):
        if not self.authenticated: return []
        try:
            r = self.session.get("https://iss.moex.com/iss/apps/nsd_corp_info/v1/accounting/msfo-full/reports", params={"companies": inn, "limit": 50})
            if r.status_code == 200 and 'json' in r.headers.get('content-type',''): return r.json().get('data', [])
        except: pass
        return []
    def get_report_details(self, report_id):
        if not self.authenticated: return {}
        try:
            r = self.session.get(f"https://iss.moex.com/iss/apps/nsd_corp_info/v1/accounting/msfo-full/reports/{report_id}", headers={"accept":"application/json"})
            if r.status_code == 200 and 'json' in r.headers.get('content-type',''): return r.json()
        except: pass
        return {}

def extract_inn(texts):
    combined = " ".join(texts)
    m = re.search(r'ИНН[:\s]*(\d{10,12})', combined)
    if m: return m.group(1)
    m = re.search(r'(?:инн|налогоплательщик|tin_code)[^\d]{0,30}(\d{10,12})', combined, re.IGNORECASE)
    if m: return m.group(1)
    return ""

def _nsd_val(vals, code):
    for v in vals:
        if v.get('line_code') == code:
            raw = v.get('value', 0); return round(float(raw)/1_000_000, 1) if raw else 0.0
    return 0.0

def nsd_to_analysis(reports_data, company_name):
    yearly = {}
    for rd in reports_data:
        if isinstance(rd, list): rd = rd[0] if rd else {}
        data = rd.get('data', rd) if isinstance(rd, dict) else {}
        values = data.get('cci_report_values', []) if isinstance(data, dict) else []
        if not values: continue
        year = None
        for v in values:
            ym = re.search(r'(\d{4})', v.get('parameter_period_name_short_ru', ''))
            if ym: year = ym.group(1); break
        if not year:
            rep = data.get('cci_report', {})
            ym = re.search(r'(\d{4})', rep.get('period_code', ''))
            if ym: year = ym.group(1)
        if year: yearly[year] = values
    if not yearly: return {}
    years = sorted(yearly.keys())
    def _sdiv(a, b): return round(a/b, 2) if b else 0.0
    def _bld(items_map, total_codes, yv):
        items, acc = [], 0.0
        for code, name in items_map.items():
            v = _nsd_val(yv, code)
            if v != 0: items.append({"name": name, "numeric_value": v}); acc += v
        total = sum(_nsd_val(yv, c) for c in total_codes)
        other = total - acc
        if abs(other) > 0.1: items.append({"name": "Прочие", "numeric_value": round(other, 1)})
        return items
    a_map = {'11111':'Денежные средства','11120':'Дебиторская задолженность','11130':'Запасы','11140':'Авансы выданные','11210':'Основные средства','11220':'Гудвил','11230':'НМА'}
    l_map = {'12110':'Кредиторская задолженность','12140':'Краткосрочные кредиты','12150':'Краткоср. часть долгоср.','12211':'Долгосрочные кредиты','12212':'Обязательства по аренде','12220':'ОНО'}
    fr_map = {'20100':'Выручка','20200':'Себестоимость','20400':'Коммерческие и адм.','21000':'Финансовые расходы','21100':'Финансовые доходы','20600':'Прочие операционные'}
    c_map = {'12310':'Уставной капитал','12320':'Добавочный капитал','12330':'Нераспределённая прибыль','12340':'Собственные акции','12370':'Неконтролирующие доли'}
    assets_comp = [{"year": y, "items": _bld(a_map, ['11000'], yearly[y])} for y in years]
    liab_comp = [{"year": y, "items": _bld(l_map, ['12100','12200'], yearly[y])} for y in years]
    fr_comp, ni_by_year = [], []
    for y in years:
        items = [{"name":n,"numeric_value":_nsd_val(yearly[y],c)} for c,n in fr_map.items() if _nsd_val(yearly[y],c)!=0]
        fr_comp.append({"year":y,"items":items}); ni_by_year.append({"year":y,"numeric_value":_nsd_val(yearly[y],'21700')})
    cap_comp = [{"year": y, "items": _bld(c_map, ['12300'], yearly[y])} for y in years]
    te_by_year = [{"year": y, "numeric_value": _nsd_val(yearly[y], '12300')} for y in years]
    ratio_defs = [
        ('Совокупный долг (млрд руб.)', lambda v: round((_nsd_val(v,'12140')+_nsd_val(v,'12150')+_nsd_val(v,'12211'))/1000, 2)),
        ('Чистый долг (млрд руб.)', lambda v: round((_nsd_val(v,'12140')+_nsd_val(v,'12150')+_nsd_val(v,'12211')-_nsd_val(v,'11111'))/1000, 2)),
        ('EBITDA (млрд руб.)', lambda v: round((_nsd_val(v,'20900')+abs(_nsd_val(v,'20800')))/1000, 2)),
        ('Рентабельность по EBITDA (%)', lambda v: _sdiv((_nsd_val(v,'20900')+abs(_nsd_val(v,'20800')))*100, _nsd_val(v,'20100'))),
        ('Чистый долг/EBITDA (x)', lambda v: _sdiv(_nsd_val(v,'12140')+_nsd_val(v,'12150')+_nsd_val(v,'12211')-_nsd_val(v,'11111'), _nsd_val(v,'20900')+abs(_nsd_val(v,'20800')))),
        ('Коэффициент автономии', lambda v: _sdiv(_nsd_val(v,'12300'), _nsd_val(v,'11000'))),
        ('ROE (%)', lambda v: _sdiv(_nsd_val(v,'21700')*100, _nsd_val(v,'12300'))),
        ('ROA (%)', lambda v: _sdiv(_nsd_val(v,'21700')*100, _nsd_val(v,'11000'))),
        ('Рентабельность по ЧП (%)', lambda v: _sdiv(_nsd_val(v,'21700')*100, _nsd_val(v,'20100'))),
        ('Процентные расходы/EBITDA (%)', lambda v: _sdiv((abs(_nsd_val(v,'21010'))+abs(_nsd_val(v,'21020')))*100, _nsd_val(v,'20900')+abs(_nsd_val(v,'20800')))),
    ]
    metrics = []
    for name, fn in ratio_defs:
        vals = [fn(yearly[y]) for y in years]
        change = f"{'+' if len(vals)>=2 and vals[-1]-vals[-2]>=0 else ''}{vals[-1]-vals[-2]:.2f}" if len(vals)>=2 else "—"
        metrics.append({"name": name, "values": vals, "change": change})
    def _chg(label, code, last, prev):
        c, p = _nsd_val(yearly[last], code), _nsd_val(yearly[prev], code)
        d = c - p; pct = _sdiv(d*100, abs(p))
        return f"{label} {'увеличились' if d > 0 else 'снизились'} на {abs(d):,.0f} млн руб. ({'+' if d > 0 else ''}{pct:.1f}%)"
    ly, py = years[-1], years[-2] if len(years) >= 2 else years[-1]
    return {
        "company_name": company_name, "reporting_period": f"{years[0]}-{years[-1]} гг." if len(years)>1 else f"{years[0]} г.",
        "general_info": {"description":"","ratings":[],"disclosure":"","beneficiaries":"","other_factors":""},
        "industry_context": {"industry_overview":"","company_position":""},
        "assets": {"total_change": _chg("Общая величина активов","11000",ly,py), "main_changes":[], "dominant_items":"", "commentary":"", "composition_by_year": assets_comp},
        "liabilities": {"total_change": _chg("Общая величина обязательств","12100",ly,py), "main_changes":[], "dominant_items":"", "commentary":"", "composition_by_year": liab_comp},
        "financial_results": {"revenue_change": _chg("Выручка","20100",ly,py), "opex_change":"", "net_income_change": _chg("Чистая прибыль","21700",ly,py), "commentary":"", "composition_by_year": fr_comp, "net_income_by_year": ni_by_year},
        "capital": {"total_change": _chg("Капитал","12300",ly,py), "commentary":"", "composition_by_year": cap_comp, "total_equity_by_year": te_by_year},
        "ratio_analysis": {"years": years, "metrics": metrics, "commentary":""},
        "conclusions": {"positives":[], "negatives":[], "credit_quality":""},
    }

# ---------------------------------------------------------------------------
# Prompts
# ---------------------------------------------------------------------------

ANALYSIS_PROMPT = """Ты — старший кредитный аналитик. Язык ответа: РУССКИЙ.
КРИТИЧЕСКИ ВАЖНО О СТРУКТУРЕ ВХОДНЫХ ДАННЫХ:
В тексте МСФО ты увидишь маркеры вида "=== НАЧАЛО ГОДОВОГО ОТЧЁТА ЗА 2025 ГОД ===" — это означает, что далее идёт годовой отчёт за указанный год. Таких блоков несколько (обычно 3-4). Каждый годовой отчёт содержит балансовый отчёт, ОПУ, ОДДС с ДВУМЯ колонками: текущий год и сравнительный предыдущий.
КАК СОБИРАТЬ ДАННЫЕ:
- Из блока "ГОДОВОЙ ОТЧЁТ ЗА 2025" бери данные за 2025 (первая колонка) и 2024 (вторая)
- Из блока "ГОДОВОЙ ОТЧЁТ ЗА 2024" бери данные за 2023 (вторая колонка)
- И так далее. РЕЗУЛЬТАТ: временной ряд за ВСЕ годы
- НЕ ОСТАНАВЛИВАЙСЯ на первом блоке! Пройди ВСЕ.
Тебе предоставлены два блока документов:
1) ОБЩИЕ ДОКУМЕНТЫ 2) ОТЧЁТНОСТЬ МСФО
Верни ТОЛЬКО валидный JSON — без markdown.
КРИТИЧНО: numeric_value — число (float). Денежные — в млн руб. Проценты — как число (42.3).
JSON структура:
"company_name": строка,
"reporting_period": строка (например "2022-2025 гг."),
"general_info": {"description": строка, "ratings": [{"agency":"","level":"","date":"","action":""}], "disclosure": строка, "beneficiaries": строка, "other_factors": строка},
"industry_context": {"industry_overview": строка (2 абзаца), "company_position": строка (2 абзаца)},
"assets": {"total_change": строка, "main_changes": [{"item":"","change_text":""}], "dominant_items": строка, "commentary": строка (CFA 1-2 абз.), "composition_by_year": [{"year":"","items":[{"name":"","numeric_value":0}]}]},
"liabilities": аналогично assets,
ПРАВИЛА composition_by_year: ВКЛЮЧИ ВСЕ ГОДЫ, 5-10 статей >5%, остальное в "Прочие", млн руб.
"financial_results": {"revenue_change": строка, "opex_change": строка, "net_income_change": строка, "commentary": строка, "composition_by_year": [...], "net_income_by_year": [{"year":"","numeric_value":0}]},
"capital": {"total_change": строка, "commentary": строка, "composition_by_year": [...], "total_equity_by_year": [{"year":"","numeric_value":0}]},
"ratio_analysis": {"years": [...], "metrics": [{"name":"","values":[...],"change":""}], "commentary": строка},
Метрики: Совокупный долг (млрд), Чистый долг (млрд), EBITDA (млрд), Рентабельность по EBITDA (%), Чистый долг/EBITDA (x), Коэф. автономии, ROE (%), ROA (%), Рентабельность по ЧП (%), Процентные расходы/EBITDA (%).
"conclusions": {"positives": [...], "negatives": [...], "credit_quality": строка}
=== ОБЩИЕ ДОКУМЕНТЫ ===
PLACEHOLDER_GENERAL_DOCS
=== ОТЧЁТНОСТЬ МСФО ===
PLACEHOLDER_IFRS_DOCS
"""

COMMENTARY_PROMPT = """Ты — старший кредитный аналитик. Язык ответа: РУССКИЙ.
Тебе предоставлены: 1) ОБЩИЕ ДОКУМЕНТЫ 2) ОТЧЁТНОСТЬ МСФО (Примечания) 3) УЖЕ ГОТОВЫЕ ФИНАНСОВЫЕ ДАННЫЕ (НЕ ПЕРЕСЧИТЫВАЙ!)
Напиши ТОЛЬКО ТЕКСТОВЫЕ КОММЕНТАРИИ. Верни валидный JSON:
{"general_info": {"description":"1 абзац","ratings":[{"agency":"","level":"","date":"","action":""}],"disclosure":"","beneficiaries":"","other_factors":""},
"industry_context": {"industry_overview":"2 абзаца","company_position":"2 абзаца"},
"assets_commentary": "CFA 1-2 абз.", "assets_dominant": "Основную долю составляют...",
"liabilities_commentary": "CFA 1-2 абз.", "liabilities_dominant": "",
"financial_results_commentary": "CFA 1-2 абз.", "capital_commentary": "CFA 1-2 абз.",
"ratio_commentary": "CFA 1-2 абз.",
"conclusions": {"positives":["3-5 факторов"],"negatives":["3-5 рисков"],"credit_quality":"2-3 предложения"}}
=== ОБЩИЕ ДОКУМЕНТЫ ===
PLACEHOLDER_GENERAL_DOCS
=== ОТЧЁТНОСТЬ МСФО (Примечания) ===
PLACEHOLDER_IFRS_DOCS
=== ФИНАНСОВЫЕ ДАННЫЕ (извлечены, НЕ ПЕРЕСЧИТЫВАЙ) ===
PLACEHOLDER_FINANCIAL_SUMMARY
"""

# ---------------------------------------------------------------------------
# LLM analysis functions
# ---------------------------------------------------------------------------

def _call_gemini(prompt, user_api_key="", max_tokens=32768):
    keys = _get_api_keys(user_api_key)
    last_error = None
    for i, key in enumerate(keys):
        try:
            client = OpenAI(api_key=key, base_url="https://openrouter.ai/api/v1")
            import time as _t; t0 = _t.time()
            r = client.chat.completions.create(
                model=_LLM_MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.2,
                max_tokens=max_tokens,
            )
            text = r.choices[0].message.content
            print(f"  LLM ответ за {_t.time()-t0:.0f} сек, {len(text):,} символов", flush=True)
            raw = text.strip()
            if raw.startswith("```"): raw = raw.split("\n",1)[1]
            if raw.endswith("```"): raw = raw.rsplit("```",1)[0]
            return json.loads(raw.strip())
        except json.JSONDecodeError:
            rep = raw.strip()
            if rep.count('"')%2!=0: rep+='"'
            rep += "]"*max(0, rep.count("[")-rep.count("]"))
            rep += "}"*max(0, rep.count("{")-rep.count("}"))
            return json.loads(rep)
        except Exception as e:
            last_error = e
            if any(w in str(e).lower() for w in ["429","quota","rate","limit"]) and i<len(keys)-1: continue
            raise
    raise last_error

def analyze_with_gemini(model, general_texts, ifrs_texts, user_api_key=""):
    gc = "\n\n".join(general_texts) if general_texts else "[Не предоставлены]"
    ic = "\n\n".join(ifrs_texts) if ifrs_texts else "[Не предоставлены]"
    if len(gc) > 150_000: gc = gc[:150_000] + "\n[...обрезано...]"
    if len(ic) > 800_000: ic = ic[:400_000] + "\n\n[...пропущено...]\n\n" + ic[-400_000:]
    found_years = sorted(set(y for y in re.findall(r'\b(20[1-3]\d)\b', ic) if 2015<=int(y)<=2035))
    if len(found_years) >= 2:
        ic = f"[СИСТЕМНАЯ ИНФОРМАЦИЯ: годы: {', '.join(found_years)}. Включи ВСЕ в composition_by_year.]\n\n" + ic
    prompt = ANALYSIS_PROMPT.replace("PLACEHOLDER_GENERAL_DOCS", gc).replace("PLACEHOLDER_IFRS_DOCS", ic)
    print(f"Размер промпта: {len(prompt):,} символов", flush=True)
    return _call_gemini(prompt, user_api_key, 32768)

def _financial_summary(analysis):
    lines = [f"Компания: {analysis['company_name']}", f"Период: {analysis['reporting_period']}",
             f"Активы: {analysis['assets']['total_change']}", f"Обязательства: {analysis['liabilities']['total_change']}",
             f"Выручка: {analysis['financial_results']['revenue_change']}", f"ЧП: {analysis['financial_results']['net_income_change']}",
             f"Капитал: {analysis['capital']['total_change']}"]
    for m in analysis.get('ratio_analysis',{}).get('metrics',[]):
        vs = ", ".join(f"{y}:{v}" for y,v in zip(analysis['ratio_analysis']['years'], m['values']))
        lines.append(f"{m['name']}: {vs} (изм. {m.get('change','—')})")
    return "\n".join(lines)

def analyze_commentary(general_texts, ifrs_texts, fin_summary, user_api_key=""):
    gc = "\n\n".join(general_texts) if general_texts else "[Не предоставлены]"
    ic = "\n\n".join(ifrs_texts) if ifrs_texts else "[Не предоставлены]"
    if len(gc)>150_000: gc=gc[:150_000]+"\n[...обрезано...]"
    if len(ic)>400_000: ic=ic[:200_000]+"\n\n[...обрезано...]\n\n"+ic[-200_000:]
    prompt = COMMENTARY_PROMPT.replace("PLACEHOLDER_GENERAL_DOCS",gc).replace("PLACEHOLDER_IFRS_DOCS",ic).replace("PLACEHOLDER_FINANCIAL_SUMMARY",fin_summary)
    print(f"Размер промпта (комментарии): {len(prompt):,} символов", flush=True)
    return _call_gemini(prompt, user_api_key, 16384)

def merge_commentary(analysis, commentary):
    if not commentary: return analysis
    if commentary.get("general_info"): analysis["general_info"] = commentary["general_info"]
    if commentary.get("industry_context"): analysis["industry_context"] = commentary["industry_context"]
    for k in ["assets","liabilities"]:
        if commentary.get(f"{k}_commentary"): analysis[k]["commentary"] = commentary[f"{k}_commentary"]
        if commentary.get(f"{k}_dominant"): analysis[k]["dominant_items"] = commentary[f"{k}_dominant"]
    if commentary.get("financial_results_commentary"): analysis["financial_results"]["commentary"] = commentary["financial_results_commentary"]
    if commentary.get("capital_commentary"): analysis["capital"]["commentary"] = commentary["capital_commentary"]
    if commentary.get("ratio_commentary"): analysis["ratio_analysis"]["commentary"] = commentary["ratio_commentary"]
    if commentary.get("conclusions"): analysis["conclusions"] = commentary["conclusions"]
    return analysis

# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------

from pptx import Presentation; from pptx.util import Inches, Pt; from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN; from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE; from pptx.chart.data import CategoryChartData

_BLUE=RGBColor(0x00,0x70,0xC0); _ACCENT=RGBColor(0x44,0x72,0xC4); _DK2=RGBColor(0x44,0x54,0x6A)
_LT2=RGBColor(0xE7,0xE6,0xE6); _BLACK=RGBColor(0,0,0); _WHITE=RGBColor(0xFF,0xFF,0xFF)
_GREY=RGBColor(0xA5,0xA5,0xA5); _FONT="Segoe UI"
_CHART_COLORS=[RGBColor(0x44,0x72,0xC4),RGBColor(0xED,0x7D,0x31),RGBColor(0xA5,0xA5,0xA5),RGBColor(0xFF,0xC0,0x00),RGBColor(0x5B,0x9B,0xD5),RGBColor(0x70,0xAD,0x47),RGBColor(0xB8,0x50,0x42),RGBColor(0x6D,0x2E,0x46),RGBColor(0x02,0x80,0x90),RGBColor(0x95,0x4F,0x72)]

def _num(v):
    if isinstance(v,(int,float)): return float(v)
    if isinstance(v,str):
        m=re.search(r"-?[\d.]+",v.replace(",",".")); return float(m.group()) if m else 0.0
    return 0.0

def _txt(sl,text,x,y,w,h,sz=10,bold=False,color=_BLACK,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=str(text); p.font.size=Pt(sz); p.font.bold=bold; p.font.color.rgb=color; p.font.name=_FONT; p.alignment=align; p.space_after=Pt(0)
    return tb

def _mtxt(sl,text,x,y,w,h,sz=10,color=_BLACK):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(str(text).split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=line; p.font.size=Pt(sz); p.font.color.rgb=color; p.font.name=_FONT; p.space_after=Pt(2)
    return tb

def _logo(sl,path):
    if path and os.path.exists(path): sl.shapes.add_picture(path,Inches(0.15),Inches(0.12),Inches(0.45),Inches(0.39))

def _header(sl,title,num,lp):
    _logo(sl,lp); _txt(sl,title,0.7,0.25,8.8,0.7,sz=24,color=_ACCENT); _txt(sl,str(num),9.2,5.2,0.5,0.3,sz=10,color=_GREY,align=PP_ALIGN.RIGHT)

def _footer(sl):
    sh=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(5.0),Inches(10),Inches(0.625)); sh.fill.solid(); sh.fill.fore_color.rgb=_BLUE; sh.line.fill.background()
    _txt(sl,"=== MVP === "*13,0.2,5.08,9.6,0.45,sz=8,color=_WHITE)

def _process_comp(cby):
    if not cby: return []
    last=cby[-1]; items=list(last.get("items",[])); total=sum(_num(it.get("numeric_value",0)) for it in items)
    if total==0: return cby
    items.sort(key=lambda it:_num(it.get("numeric_value",0)),reverse=True); sel=set()
    for it in items:
        if _num(it["numeric_value"])/total>=0.05: sel.add(it["name"])
    for it in items[:5]: sel.add(it["name"])
    sl=list(sel)[:10]; result=[]
    for yo in cby:
        kept,other=[],0.0
        for it in yo.get("items",[]):
            if it["name"] in sl: kept.append(it)
            else: other+=_num(it.get("numeric_value",0))
        if other>0: kept.append({"name":"Прочие","numeric_value":other})
        result.append({"year":yo["year"],"items":kept})
    return result

def _stacked_chart(sl,cby,x,y,w,h):
    proc=_process_comp(cby)
    if not proc: return None
    names,seen=[],set()
    for yr in proc:
        for it in yr.get("items",[]):
            if it["name"] not in seen: names.append(it["name"]); seen.add(it["name"])
    years=[yr["year"] for yr in proc]; cd=CategoryChartData(); cd.categories=years
    for name in names:
        vals=[_num(next((it for it in yr.get("items",[]) if it["name"]==name),{"numeric_value":0}).get("numeric_value",0)) for yr in proc]
        cd.add_series(name,vals)
    cf=sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,Inches(x),Inches(y),Inches(w),Inches(h),cd)
    ch=cf.chart; ch.has_legend=True; ch.legend.include_in_layout=False; ch.legend.font.size=Pt(7); ch.legend.font.name=_FONT
    ch.plots[0].gap_width=60
    for i,s in enumerate(ch.series):
        s.format.fill.solid(); s.format.fill.fore_color.rgb=_CHART_COLORS[i%len(_CHART_COLORS)]; s.invert_if_negative=False
        s.has_data_labels=True; dl=s.data_labels; dl.font.size=Pt(7); dl.font.color.rgb=_WHITE; dl.font.name=_FONT; dl.number_format='#,##0'; dl.show_value=True
    ch.category_axis.tick_labels.font.size=Pt(9); ch.category_axis.tick_labels.font.name=_FONT; ch.category_axis.tick_labels.font.color.rgb=_DK2
    ch.value_axis.tick_labels.font.size=Pt(8); ch.value_axis.tick_labels.font.name=_FONT; ch.value_axis.tick_labels.font.color.rgb=_DK2
    ch.value_axis.major_gridlines.format.line.color.rgb=_LT2
    return cf

def _add_combo_line(chart_frame, line_data, series_name, line_color_hex, label_color=None):
    if not chart_frame or not line_data: return
    from lxml import etree
    C='http://schemas.openxmlformats.org/drawingml/2006/chart'; A='http://schemas.openxmlformats.org/drawingml/2006/main'
    pa=chart_frame.chart.element.find(f'{{{C}}}chart/{{{C}}}plotArea')
    if pa is None: return
    ca=pa.find(f'{{{C}}}catAx'); bc=pa.find(f'{{{C}}}barChart')
    if ca is None or bc is None: return
    cid=ca.find(f'{{{C}}}axId').get('val'); sid=str(int(cid)+100); ni=len(bc.findall(f'{{{C}}}ser'))
    lc=etree.SubElement(pa,f'{{{C}}}lineChart'); etree.SubElement(lc,f'{{{C}}}grouping').set('val','standard'); etree.SubElement(lc,f'{{{C}}}varyColors').set('val','0')
    ser=etree.SubElement(lc,f'{{{C}}}ser'); etree.SubElement(ser,f'{{{C}}}idx').set('val',str(ni)); etree.SubElement(ser,f'{{{C}}}order').set('val',str(ni))
    tx=etree.SubElement(ser,f'{{{C}}}tx'); sc=etree.SubElement(etree.SubElement(tx,f'{{{C}}}strRef'),f'{{{C}}}strCache')
    etree.SubElement(sc,f'{{{C}}}ptCount').set('val','1'); p0=etree.SubElement(sc,f'{{{C}}}pt'); p0.set('idx','0'); etree.SubElement(p0,f'{{{C}}}v').text=series_name
    sp=etree.SubElement(ser,f'{{{C}}}spPr'); ln=etree.SubElement(sp,f'{{{A}}}ln'); ln.set('w','25400')
    sf=etree.SubElement(ln,f'{{{A}}}solidFill'); etree.SubElement(sf,f'{{{A}}}srgbClr').set('val',line_color_hex)
    mk=etree.SubElement(ser,f'{{{C}}}marker'); etree.SubElement(mk,f'{{{C}}}symbol').set('val','circle'); etree.SubElement(mk,f'{{{C}}}size').set('val','5')
    ms=etree.SubElement(mk,f'{{{C}}}spPr'); mf=etree.SubElement(ms,f'{{{A}}}solidFill'); etree.SubElement(mf,f'{{{A}}}srgbClr').set('val','FFFFFF')
    ml=etree.SubElement(ms,f'{{{A}}}ln'); ml.set('w','12700'); mls=etree.SubElement(ml,f'{{{A}}}solidFill'); etree.SubElement(mls,f'{{{A}}}srgbClr').set('val',line_color_hex)
    dl=etree.SubElement(ser,f'{{{C}}}dLbls'); tp=etree.SubElement(dl,f'{{{C}}}txPr'); etree.SubElement(tp,f'{{{A}}}bodyPr'); etree.SubElement(tp,f'{{{A}}}lstStyle')
    dp=etree.SubElement(tp,f'{{{A}}}p'); ppr=etree.SubElement(dp,f'{{{A}}}pPr'); rpr=etree.SubElement(ppr,f'{{{A}}}defRPr'); rpr.set('sz','800'); rpr.set('b','1')
    rf=etree.SubElement(rpr,f'{{{A}}}solidFill'); etree.SubElement(rf,f'{{{A}}}srgbClr').set('val',label_color or line_color_hex)
    etree.SubElement(rpr,f'{{{A}}}latin').set('typeface',_FONT); etree.SubElement(dp,f'{{{A}}}endParaRPr')
    nf=etree.SubElement(dl,f'{{{C}}}numFmt'); nf.set('formatCode','#,##0'); nf.set('sourceLinked','0')
    etree.SubElement(dl,f'{{{C}}}showLegendKey').set('val','0'); etree.SubElement(dl,f'{{{C}}}showVal').set('val','1')
    etree.SubElement(dl,f'{{{C}}}showCatName').set('val','0'); etree.SubElement(dl,f'{{{C}}}showSerName').set('val','0')
    cat=etree.SubElement(ser,f'{{{C}}}cat'); sc2=etree.SubElement(etree.SubElement(cat,f'{{{C}}}strRef'),f'{{{C}}}strCache')
    etree.SubElement(sc2,f'{{{C}}}ptCount').set('val',str(len(line_data)))
    for j,item in enumerate(line_data): p2=etree.SubElement(sc2,f'{{{C}}}pt'); p2.set('idx',str(j)); etree.SubElement(p2,f'{{{C}}}v').text=str(item.get("year",""))
    v=etree.SubElement(ser,f'{{{C}}}val'); nc=etree.SubElement(etree.SubElement(v,f'{{{C}}}numRef'),f'{{{C}}}numCache')
    etree.SubElement(nc,f'{{{C}}}formatCode').text='#,##0'; etree.SubElement(nc,f'{{{C}}}ptCount').set('val',str(len(line_data)))
    for j,item in enumerate(line_data): p3=etree.SubElement(nc,f'{{{C}}}pt'); p3.set('idx',str(j)); etree.SubElement(p3,f'{{{C}}}v').text=str(_num(item.get("numeric_value",0)))
    etree.SubElement(lc,f'{{{C}}}smooth').set('val','0'); etree.SubElement(lc,f'{{{C}}}axId').set('val',cid); etree.SubElement(lc,f'{{{C}}}axId').set('val',sid)
    sa=etree.SubElement(pa,f'{{{C}}}valAx'); etree.SubElement(sa,f'{{{C}}}axId').set('val',sid)
    s3=etree.SubElement(sa,f'{{{C}}}scaling'); etree.SubElement(s3,f'{{{C}}}orientation').set('val','minMax')
    etree.SubElement(sa,f'{{{C}}}delete').set('val','1'); etree.SubElement(sa,f'{{{C}}}axPos').set('val','r'); etree.SubElement(sa,f'{{{C}}}crossAx').set('val',cid)

def generate_pptx(analysis, output_dir):
    logo_path=os.path.join(os.path.dirname(os.path.abspath(__file__)),"logo.png")
    if not os.path.exists(logo_path): logo_path=None
    cn=analysis.get("company_name","<КОМПАНИЯ>"); prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625); bl=prs.slide_layouts[6]
    # S1 Title
    s1=prs.slides.add_slide(bl); _logo(s1,logo_path); _txt(s1,cn,1,2.0,8,1.0,sz=32,align=PP_ALIGN.CENTER); _txt(s1,"КРЕДИТНОЕ ЗАКЛЮЧЕНИЕ",1,2.9,8,0.5,sz=16,align=PP_ALIGN.CENTER); _footer(s1); _txt(s1,"1",9.2,5.2,0.5,0.3,sz=10,color=_GREY,align=PP_ALIGN.RIGHT)
    # S2 General info
    s2=prs.slides.add_slide(bl); _header(s2,"ОБЩАЯ ИНФОРМАЦИЯ",2,logo_path); gi=analysis.get("general_info",{})
    _mtxt(s2,gi.get("description","Информация не предоставлена."),0.4,1.0,9.2,0.65,sz=11)
    ry=1.7; ratings=gi.get("ratings",[])
    if ratings:
        hdr=["Агентство","Рейтинг","Дата","Действие"]; rows=[[r.get("agency",""),r.get("level",""),r.get("date",""),r.get("action","")] for r in ratings]; all_r=[hdr]+rows; nr=len(all_r)
        ts=s2.shapes.add_table(nr,4,Inches(1.9),Inches(ry),Inches(7.6),Inches(0.28*nr)); tbl=ts.table
        for ci,w in enumerate([Inches(2.2),Inches(2.2),Inches(1.2),Inches(2.0)]): tbl.columns[ci].width=w
        for ri,row in enumerate(all_r):
            for ci,txt in enumerate(row):
                c=tbl.cell(ri,ci); c.text=txt; p=c.text_frame.paragraphs[0]; p.font.size=Pt(8); p.font.name=_FONT
                if ri==0: p.font.bold=True; p.font.color.rgb=_WHITE; c.fill.solid(); c.fill.fore_color.rgb=_ACCENT
                else: c.fill.solid(); c.fill.fore_color.rgb=_WHITE if ri%2==1 else _LT2
        ry+=max(0.28*nr,0.55)+0.15
    for label,text in [("Раскрытие",gi.get("disclosure","—")),("Бенефициары",gi.get("beneficiaries","—")),("Иные факторы",gi.get("other_factors","—"))]:
        sh=s2.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.4),Inches(ry),Inches(9.2),Pt(1)); sh.fill.solid(); sh.fill.fore_color.rgb=_LT2; sh.line.fill.background(); ry+=0.08
        _txt(s2,label,0.4,ry,1.4,0.55,sz=10); _txt(s2,text,1.9,ry,7.6,0.55,sz=10); ry+=0.55
    # S3 Industry
    s3=prs.slides.add_slide(bl); _header(s3,"ОТРАСЛЕВОЙ КОНТЕКСТ",3,logo_path); ic=analysis.get("industry_context",{})
    _mtxt(s3,ic.get("industry_overview",""),0.5,1.1,9,1.8,sz=11); _mtxt(s3,ic.get("company_position",""),0.5,3.0,9,1.8,sz=11)
    # S4 Assets
    s4=prs.slides.add_slide(bl); _header(s4,f"АКТИВЫ {cn}",4,logo_path); a=analysis.get("assets",{})
    _stacked_chart(s4,a.get("composition_by_year"),0.3,1.1,4.5,3.8)
    tp=[]; 
    if a.get("total_change"): tp.append(a["total_change"])
    if a.get("main_changes"):
        tp.append("Основные изменения:")
        for ch in a["main_changes"]: tp.append(f"  • {ch.get('change_text') or ch.get('item','')}")
    if a.get("dominant_items"): tp.append(""); tp.append(a["dominant_items"])
    _mtxt(s4,"\n".join(tp) or "—",5.0,1.1,4.7,2.0,sz=9)
    if a.get("commentary"): _mtxt(s4,a["commentary"],5.0,3.2,4.7,1.7,sz=9,color=_DK2)
    # S5 Liabilities
    s5=prs.slides.add_slide(bl); _header(s5,f"ОБЯЗАТЕЛЬСТВА {cn}",5,logo_path); li=analysis.get("liabilities",{})
    _stacked_chart(s5,li.get("composition_by_year"),0.3,1.1,4.5,3.8)
    tp2=[]
    if li.get("total_change"): tp2.append(li["total_change"])
    if li.get("main_changes"):
        tp2.append("Основные изменения:")
        for ch in li["main_changes"]: tp2.append(f"  • {ch.get('change_text') or ch.get('item','')}")
    if li.get("dominant_items"): tp2.append(""); tp2.append(li["dominant_items"])
    _mtxt(s5,"\n".join(tp2) or "—",5.0,1.1,4.7,2.0,sz=9)
    if li.get("commentary"): _mtxt(s5,li["commentary"],5.0,3.2,4.7,1.7,sz=9,color=_DK2)
    # S6 Financial results (combo chart)
    s6=prs.slides.add_slide(bl); _header(s6,f"ФИНАНСОВЫЕ РЕЗУЛЬТАТЫ {cn}",6,logo_path); fr=analysis.get("financial_results",{})
    cf6=_stacked_chart(s6,fr.get("composition_by_year"),0.3,1.1,4.5,3.8)
    ni=fr.get("net_income_by_year",[])
    if ni: _add_combo_line(cf6,ni,"Чистая прибыль","ED7D31","000000")
    fp=[]
    if fr.get("revenue_change"): fp.append(fr["revenue_change"])
    if fr.get("opex_change"): fp.append(fr["opex_change"])
    if fr.get("net_income_change"): fp.append(fr["net_income_change"])
    _mtxt(s6,"\n\n".join(fp) or "—",5.0,1.1,4.7,2.0,sz=9)
    if fr.get("commentary"): _mtxt(s6,fr["commentary"],5.0,3.2,4.7,1.7,sz=9,color=_DK2)
    # S7 Capital (combo chart)
    s7=prs.slides.add_slide(bl); _header(s7,f"КАПИТАЛ {cn}",7,logo_path); cap=analysis.get("capital",{})
    cf7=_stacked_chart(s7,cap.get("composition_by_year"),0.3,1.1,4.5,3.8)
    te=cap.get("total_equity_by_year",[])
    if te: _add_combo_line(cf7,te,"Итого капитал","000000","000000")
    if cap.get("total_change"): _mtxt(s7,cap["total_change"],5.0,1.1,4.7,1.5,sz=9)
    if cap.get("commentary"): _mtxt(s7,cap["commentary"],5.0,2.7,4.7,2.2,sz=9,color=_DK2)
    # S8 Ratios table
    s8=prs.slides.add_slide(bl); _header(s8,f"КОЭФФИЦИЕНТНЫЙ АНАЛИЗ {cn}",8,logo_path); ra=analysis.get("ratio_analysis",{})
    ra_y=ra.get("years",[]); ra_m=ra.get("metrics",[]); HB=RGBColor(0x30,0x54,0x96)
    if ra_y and ra_m:
        cl=["Показатель"]+ra_y+["Изм."]; nc=len(cl); nr=1+len(ra_m)
        ts=s8.shapes.add_table(nr,nc,Inches(0.3),Inches(1.1),Inches(9.4),Inches(0.3*nr)); tbl=ts.table
        tbl.columns[0].width=Inches(3.0)
        for ci in range(1,nc): tbl.columns[ci].width=Inches(6.4/(nc-1))
        for ci,lb in enumerate(cl):
            c=tbl.cell(0,ci); c.text=lb; p=c.text_frame.paragraphs[0]; p.font.size=Pt(8); p.font.bold=True; p.font.color.rgb=_WHITE; p.font.name=_FONT; p.alignment=PP_ALIGN.CENTER; c.fill.solid(); c.fill.fore_color.rgb=HB
        for ri,met in enumerate(ra_m):
            bg=_WHITE if ri%2==0 else _LT2
            c=tbl.cell(ri+1,0); c.text=met.get("name",""); p=c.text_frame.paragraphs[0]; p.font.size=Pt(7); p.font.name=_FONT; p.font.bold=True; c.fill.solid(); c.fill.fore_color.rgb=bg
            for vi,val in enumerate(met.get("values",[])):
                if vi+1<nc-1:
                    c=tbl.cell(ri+1,vi+1); c.text=str(round(val,2)) if isinstance(val,float) else str(val); p=c.text_frame.paragraphs[0]; p.font.size=Pt(7); p.font.name=_FONT; p.alignment=PP_ALIGN.CENTER; c.fill.solid(); c.fill.fore_color.rgb=bg
            c=tbl.cell(ri+1,nc-1); c.text=met.get("change","—"); p=c.text_frame.paragraphs[0]; p.font.size=Pt(7); p.font.name=_FONT; p.font.bold=True; p.alignment=PP_ALIGN.CENTER; c.fill.solid(); c.fill.fore_color.rgb=bg
    if ra.get("commentary"): _mtxt(s8,ra["commentary"],0.3,1.1+0.3*(len(ra_m)+1)+0.2,9.4,1.8,sz=9,color=_DK2)
    # S9 Conclusions
    s9=prs.slides.add_slide(bl); _header(s9,"ИТОГИ",9,logo_path); co=analysis.get("conclusions",{})
    GF=RGBColor(0x70,0xAD,0x47); NF=RGBColor(0xF4,0xB9,0xA4); NO=RGBColor(0xED,0x7D,0x31)
    sh=s9.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.3),Inches(1.0),Inches(4.4),Inches(0.35)); sh.fill.solid(); sh.fill.fore_color.rgb=GF; sh.line.color.rgb=GF
    _txt(s9,"ПОЛОЖИТЕЛЬНЫЕ СТОРОНЫ",0.4,1.02,4.2,0.3,sz=10,bold=True,color=_WHITE)
    sh=s9.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.3),Inches(1.4),Inches(4.4),Inches(2.0)); sh.fill.solid(); sh.fill.fore_color.rgb=_WHITE; sh.line.color.rgb=GF; sh.line.width=Pt(1.5)
    _mtxt(s9,"\n".join(f"• {p}" for p in co.get("positives",[])) or "—",0.45,1.5,4.1,1.8,sz=9)
    sh=s9.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(5.2),Inches(1.0),Inches(4.4),Inches(0.35)); sh.fill.solid(); sh.fill.fore_color.rgb=NF; sh.line.color.rgb=NO
    _txt(s9,"НЕГАТИВНЫЕ СТОРОНЫ",5.3,1.02,4.2,0.3,sz=10,bold=True,color=_WHITE)
    sh=s9.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(5.2),Inches(1.4),Inches(4.4),Inches(2.0)); sh.fill.solid(); sh.fill.fore_color.rgb=_WHITE; sh.line.color.rgb=NO; sh.line.width=Pt(1.5)
    _mtxt(s9,"\n".join(f"• {n}" for n in co.get("negatives",[])) or "—",5.35,1.5,4.1,1.8,sz=9)
    if co.get("credit_quality"): _mtxt(s9,co["credit_quality"],0.3,3.8,9.4,1.0,sz=10,color=_DK2)
    company=(cn).replace(" ","_").replace("/","_")[:40]; out=os.path.join(output_dir,f"{company}_Кредитное_Заключение.pptx"); prs.save(out); return out

# ---------------------------------------------------------------------------
# Pipelines
# ---------------------------------------------------------------------------

def process_standard(api_key, general_files, ifrs_files, progress=gr.Progress()):
    """Стандарт: финансы из НРД + текст из LLM."""
    if not general_files and not ifrs_files: raise gr.Error("Загрузите хотя бы один документ.")
    progress(0.05, desc="Извлечение текста..."); general_texts=extract_texts(general_files); ifrs_texts=extract_texts(ifrs_files)
    if not general_texts and not ifrs_texts: raise gr.Error("Не удалось извлечь текст.")
    progress(0.10, desc="Поиск ИНН..."); inn=extract_inn(general_texts+ifrs_texts)
    if not inn: raise gr.Error("ИНН не найден в документах.")
    print(f"  ИНН: {inn}", flush=True)
    progress(0.15, desc="Подключение к MOEX НРД..."); client=MOEXClient()
    if not client.login(): raise gr.Error("Ошибка авторизации НРД. Проверьте секрет 'NSD'.")
    progress(0.20, desc="Запрос отчётов из НРД..."); rl=client.get_report_list(inn)
    if not rl: raise gr.Error(f"Отчёты для ИНН {inn} не найдены.")
    cn=rl[0].get('company_name_short_ru','<Компания>'); print(f"  {cn}, {len(rl)} отчётов", flush=True)
    progress(0.30, desc="Загрузка данных..."); rds=[]
    for rep in rl[:5]:
        rid=rep.get('basis_type_report_id')
        if rid: d=client.get_report_details(rid); (rds.append(d) if d else None)
    if not rds: raise gr.Error("Не удалось загрузить отчёты.")
    progress(0.45, desc="Расчёт показателей (детерминированный)..."); analysis=nsd_to_analysis(rds,cn)
    if not analysis: raise gr.Error("Ошибка преобразования данных.")
    progress(0.55, desc="Запрос комментариев у LLM...")
    try:
        fs=_financial_summary(analysis); commentary=analyze_commentary(general_texts,ifrs_texts,fs,api_key or ""); analysis=merge_commentary(analysis,commentary)
    except Exception as e: print(f"⚠️ Комментарии: {e}", flush=True)
    progress(0.80, desc="Генерация PPTX..."); od=tempfile.mkdtemp()
    try: pp=generate_pptx(analysis,od); sp=os.path.join(tempfile.gettempdir(),os.path.basename(pp)); shutil.copy2(pp,sp)
    except Exception as e: raise gr.Error(f"Ошибка PPTX: {e}")
    finally: shutil.rmtree(od,ignore_errors=True)
    progress(0.95, desc="Готово!")
    return sp, f"## {cn} — {analysis.get('reporting_period','')}\n\n**Тариф:** Стандарт (НРД, ИНН: {inn})\n\n**Активы:** {analysis['assets']['total_change']}\n\n**Выручка:** {analysis['financial_results']['revenue_change']}"

def process_premium(api_key, general_files, ifrs_files, progress=gr.Progress()):
    """Премиум: всё через LLM."""
    if not general_files and not ifrs_files: raise gr.Error("Загрузите хотя бы один документ.")
    progress(0.05, desc="Настройка LLM..."); model=get_gemini_model(api_key or "")
    progress(0.10, desc="Извлечение текста..."); gt=extract_texts(general_files)
    progress(0.15, desc="Извлечение МСФО (OCR ~1-2 мин)..."); it=extract_texts(ifrs_files)
    if not gt and not it: raise gr.Error("Не удалось извлечь текст.")
    progress(0.35, desc="Анализ LLM (~30-60 сек)...")
    try: analysis=analyze_with_gemini(model,gt,it,api_key or "")
    except json.JSONDecodeError as e: raise gr.Error(f"Невалидный JSON: {e}")
    except Exception as e: raise gr.Error(f"Ошибка: {e}")
    progress(0.7, desc="Генерация PPTX..."); od=tempfile.mkdtemp()
    try: pp=generate_pptx(analysis,od); sp=os.path.join(tempfile.gettempdir(),os.path.basename(pp)); shutil.copy2(pp,sp)
    except Exception as e: raise gr.Error(f"Ошибка PPTX: {e}")
    finally: shutil.rmtree(od,ignore_errors=True)
    progress(0.95, desc="Готово!")
    cn=analysis.get("company_name","Компания"); rp=analysis.get("reporting_period","")
    sm=f"## {cn} — {rp}\n\n**Тариф:** Премиум (полный LLM-анализ)\n\n"
    if analysis.get("assets",{}).get("total_change"): sm+=f"**Активы:** {analysis['assets']['total_change']}\n\n"
    return sp, sm

# ---------------------------------------------------------------------------
# Gradio UI
# ---------------------------------------------------------------------------

has_env_key=bool(os.environ.get("OPENROUTER_KEY","").strip())
CSS="""
.gradio-container{max-width:960px!important}
.main-header{text-align:center;margin-bottom:0.5em}
.main-header h1{font-size:1.8em;font-weight:700;color:#e0e0e0}
.main-header p{color:#94a3b8;font-size:0.95em}
.env-hint{padding:0.5em 1em;background:rgba(255,255,255,0.05);border-left:3px solid #4472C4;border-radius:4px;margin-bottom:1em;font-size:0.9em;color:#b0b0b0}
.reset-btn{background:#dc2626!important;color:white!important;border:none!important;font-weight:700!important}
.standard-btn{background:#2E7D32!important;color:white!important;font-weight:700!important}
.premium-btn{background:#4472C4!important;color:white!important;font-weight:700!important}
footer{display:none!important}
"""

_logo_html=""
_logo_path_ui=os.path.join(os.path.dirname(os.path.abspath(__file__)),"logo.png")
if os.path.exists(_logo_path_ui):
    import base64 as _b64
    with open(_logo_path_ui,"rb") as _lf: _lb=_lf.read()
    _logo_html=f'<img src="data:image/png;base64,{_b64.b64encode(_lb).decode()}" style="height:50px;margin-bottom:0.3em"/>'

with gr.Blocks(title="Генератор Кредитных Отчетов") as demo:
    gr.HTML(f'<div class="main-header">{_logo_html}<h1>ГЕНЕРАТОР КРЕДИТНЫХ ОТЧЕТОВ (MVP)</h1><p>Загрузите документы — получите кредитное заключение</p></div>')
    if has_env_key:
        gr.HTML('<div class="env-hint">Инфраструктура настроена.</div>'); api_key_input=gr.Textbox(value="",visible=False)
    else:
        api_key_input=gr.Textbox(label="OpenRouter API Key",placeholder="sk-or-...",type="password")
    with gr.Row():
        with gr.Column(): general_input=gr.File(label="Общие документы",file_count="multiple",file_types=[".pdf",".txt",".docx"],type="filepath")
        with gr.Column(): ifrs_input=gr.File(label="Отчётность МСФО",file_count="multiple",file_types=[".pdf",".txt",".docx"],type="filepath")
    with gr.Row():
        standard_btn=gr.Button('Тариф "Стандарт" (НРД + LLM)',variant="primary",size="lg",scale=2,elem_classes=["standard-btn"])
        premium_btn=gr.Button('Тариф "Премиум" (полный LLM)',variant="primary",size="lg",scale=2,elem_classes=["premium-btn"])
        reset_btn=gr.Button("RESET!",size="lg",scale=1,elem_classes=["reset-btn"])
    gr.HTML('<div class="env-hint"><b>Стандарт:</b> финансовые данные из НРД (по ИНН), текст из LLM. Точные числа, быстро.<br/><b>Премиум:</b> весь анализ через LLM из загруженных PDF. Гибче, но числа могут отличаться.</div>')
    with gr.Row():
        with gr.Column(): output_file=gr.File(label="Скачать презентацию",interactive=False)
        with gr.Column(): output_summary=gr.Markdown(label="Предпросмотр")
    standard_btn.click(fn=process_standard,inputs=[api_key_input,general_input,ifrs_input],outputs=[output_file,output_summary])
    premium_btn.click(fn=process_premium,inputs=[api_key_input,general_input,ifrs_input],outputs=[output_file,output_summary])
    def reset_all(): return None,None,None,""
    reset_btn.click(fn=reset_all,inputs=[],outputs=[general_input,ifrs_input,output_file,output_summary])
    gr.HTML('<div style="text-align:center;margin-top:1em;padding:0.8em;border-top:1px solid rgba(255,255,255,0.1)"><p style="color:#64748b;font-size:0.85em">Denis Pokrovsky · OpenRouter · НРД API · python-pptx</p></div>')

if __name__=="__main__": demo.launch(server_name="0.0.0.0",server_port=7860,css=CSS,theme=gr.themes.Default(primary_hue="blue",neutral_hue="slate").set(body_background_fill="#1a1a2e",block_background_fill="#16213e",input_background_fill="#0f3460",button_primary_background_fill="#4472C4",button_primary_text_color="#ffffff"))
